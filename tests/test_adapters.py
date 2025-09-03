from docflow.adapters.docx_adapter import DocxAdapter
from docflow.adapters.pptx_adapter import PptxAdapter
import io
from docx import Document
from pptx import Presentation
from pptx.util import Inches


def test_docx_adapter_text_and_image(tmp_path):
    t = tmp_path / 'template.docx'
    doc = Document()
    doc.add_paragraph('Hello {{name}}')
    doc.add_paragraph('Image here: {{image:hero}}')
    doc.save(t)

    adapter = DocxAdapter(str(t))
    adapter.load()
    pls = adapter.list_placeholders()
    assert 'name' in pls
    assert 'hero' in pls

    # prepare real png bytes via matplotlib so python-docx recognizes it
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots()
    ax.text(0.5, 0.5, 'Alice', ha='center')
    bio = io.BytesIO()
    fig.savefig(bio, format='png')
    bio.seek(0)
    mapping = {'name': 'Alice', 'hero': bio.getvalue()}
    adapter.apply(mapping=mapping, global_vars={})
    out = tmp_path / 'out.docx'
    adapter.save(str(out))
    assert out.exists()


def test_pptx_adapter_text_and_image(tmp_path):
    t = tmp_path / 'template.pptx'
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tx.text = 'Title: {{title}}\n{{image:img}}'
    prs.save(t)

    adapter = PptxAdapter(str(t))
    adapter.load()
    pls = adapter.list_placeholders()
    assert 'title' in pls
    assert 'img' in pls

    # we create a real small png via matplotlib so pptx can accept it
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    fig, ax = plt.subplots()
    ax.text(0.5, 0.5, 'hi')
    bio = io.BytesIO()
    fig.savefig(bio, format='png')
    bio.seek(0)
    adapter.apply(mapping={'title': 'Report', 'img': bio.getvalue()}, global_vars={})
    out = tmp_path / 'out.pptx'
    adapter.save(str(out))
    assert out.exists()
