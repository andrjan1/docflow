from pathlib import Path
from docflow.cli import init as cli_init
from docflow.runtime.orchestrator import run_config


def test_examples_e2e(tmp_path):
    # ensure example config and prompts exist
    cli_init('config/example.config.yaml')

    # create kb folder and a simple doc
    kbdir = Path('kb')
    kbdir.mkdir(exist_ok=True)
    (kbdir / 'note.md').write_text('Questo è un documento di KB di esempio sul prodotto Linea X')

    # create templates
    tdir = Path('templates')
    tdir.mkdir(parents=True, exist_ok=True)
    # docx
    from docx import Document

    doc = Document()
    doc.add_paragraph('Titolo: {{abstract}}')
    doc.add_paragraph('{{image:hero_img}}')
    doc.add_paragraph('{{image:sales_chart}}')
    doc.save(tdir / 'demo_template.docx')

    # pptx
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tx.text = 'Slide: {{title}}\n{{image:cover_img}}'
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    tx2 = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tx2.text = '{{image:chart_img}}'
    prs.save(tdir / 'demo_template.pptx')

    # run orchestrator
    res = run_config('config/example.config.yaml', verbose=True)
    assert isinstance(res, dict)
    out_files = list(Path('build/output').iterdir())
    assert any(f.suffix in ('.docx', '.pptx') for f in out_files)
