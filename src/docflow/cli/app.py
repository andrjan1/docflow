import typer
from pathlib import Path
from ..config import load_config
from ..runtime.orchestrator import run_config
from ..adapters.docx_adapter import DocxAdapter
from ..adapters.pptx_adapter import PptxAdapter
from typing import List
from rich.table import Table
from rich.console import Console
from ..kb.strategies import prepare_kb_for_action
from ..runtime.prompt_builder import build_prompt_for_action
# attempt to load .env from project root for CLI runs
try:
    from dotenv import load_dotenv
    load_dotenv()
except Exception:
    pass

app = typer.Typer()


@app.command()
def config_validate(config: str):
    load_config(config)
    typer.echo('valid')


@app.command()
def inspect_template(template: str, adapter: str = 'docx'):
    adapters = {'docx': DocxAdapter, 'pptx': PptxAdapter}
    AdapterCls = adapters.get(adapter)
    if not AdapterCls:
        typer.echo('unknown adapter')
        raise typer.Exit(code=1)
    a = AdapterCls(template)
    typer.echo('\n'.join(a.list_placeholders()))


@app.command()
def dry_run(config: str, verbose: bool = False):
    cfg = load_config(config)
    # print order and deps
    def _dump(a):
        if hasattr(a, 'model_dump'):
            return a.model_dump()
        if hasattr(a, 'dict'):
            return a.dict()
        return a

    actions = [_dump(a) for a in cfg.workflow.actions]
    from ..core.workflow import toposort_actions

    order = toposort_actions(actions)
    typer.echo('Action order:')
    for a in order:
        typer.echo(f"- {a['id']} deps={a.get('deps',[])}")
        kb_text = prepare_kb_for_action(a, {})
        prompt = build_prompt_for_action(a, {}, kb_text)
        typer.echo(f"  prompt: {prompt}")
    if verbose:
        # show simple telemetry if available (dry-run uses no ctx so limited info)
        typer.echo('\nDry-run complete')


@app.command()
def run(config: str, verbose: bool = False):
    res = run_config(config, verbose=verbose)
    typer.echo(f'Wrote {len(res)} files')
    if verbose:
        console = Console()
        # show files and basic counts
        table = Table('file')
        table.add_column('file')
        for f in res.keys():
            table.add_row(f)
        console.print(table)


@app.command()
def init(name: str = 'docflow', provider: str = 'mock', with_kb: bool = False, prompt_mode: str = 'inline', adapters: List[str] = ['docx','pptx']):
    # create config file and example templates
    p = Path('config/example.config.yaml')
    p.parent.mkdir(parents=True, exist_ok=True)
    if not p.exists():
        # write an AppConfig-shaped example config so init works from any CWD (e.g. tests tmp_path)
        example = (
            "project:\n"
            "  base_dir: ..\n"
            "  output_dir: build/output\n"
            "  temp_dir: build/tmp\n"
            "  log_level: INFO\n\n"
            "ai:\n"
            "  provider: mock\n\n"
            "workflow:\n"
            "  templates:\n"
            "    - path: templates/demo_template.docx\n"
            "      adapter: docx\n"
            "    - path: templates/demo_template.pptx\n"
            "      adapter: pptx\n"
            "  actions:\n"
            "    - id: gen1\n"
            "      type: generative\n"
            "      returns: image\n"
            "      prompt: |\n"
            "        Say hello to {{name}}\n"
            "    - id: code1\n"
            "      type: code\n"
            "      returns: image\n"
            "      code: |\n"
            "        def run(vars):\n"
            "            import matplotlib.pyplot as plt\n"
            "            import io\n"
            "            fig, ax = plt.subplots()\n"
            "            ax.plot([1,2,3],[1,4,9])\n"
            "            ax.set_title(vars.get('greeting','chart'))\n"
            "            bio = io.BytesIO()\n"
            "            fig.savefig(bio, format='png')\n"
            "            bio.seek(0)\n"
            "            return {'vars': {'charted': True}, 'image_chart': bio.getvalue()}\n"
        )
        p.write_text(example)
    tdir = Path('templates')
    tdir.mkdir(parents=True, exist_ok=True)
    # create real templates using python-docx and python-pptx if available
    try:
        from docx import Document
        from pptx import Presentation
        from pptx.util import Inches
        doc = Document()
        doc.add_paragraph('Hello {{greeting}}')
        doc.add_paragraph('Name: {{name}}')
        doc.add_paragraph('Image: {{image:image}}')
        doc.save(tdir / 'demo_template.docx')
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tx = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
        tx.text = 'Slide: {{greeting}}\nName: {{name}}\n{{image:image}}'
        prs.save(tdir / 'demo_template.pptx')
        typer.echo('Created templates in templates/')
    except Exception as e:
        typer.echo('Could not create Office templates (missing libs): ' + str(e))


if __name__ == '__main__':
    app()
