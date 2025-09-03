from typing import Dict, Any, List
from docflow.adapters.base import DocumentAdapter
import re
import io
from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

IMAGE_RE = re.compile(r"{{image:([^}]+)}}")
VAR_RE = re.compile(r"{{\s*([^}]+)\s*}}")


class PptxAdapter(DocumentAdapter):
    def __init__(self, path: str):
        super().__init__(Path(path))
        self.prs = None

    def load(self):
        self.prs = Presentation(self.path)

    def list_placeholders(self) -> List[str]:
        if self.prs is None:
            self.load()
        found = set()
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                text = shape.text or ''
                for m in VAR_RE.finditer(text):
                    name = m.group(1)
                    if ':' in name:
                        name = name.split(':', 1)[1]
                    found.add(name)
        return sorted(found)

    def apply(self, mapping: Dict[str, Any], global_vars: Dict[str, Any]):
        if self.prs is None:
            self.load()
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not getattr(shape, 'has_text_frame', False):
                    continue
                text = shape.text or ''
                m = IMAGE_RE.search(text)
                if m:
                    key = m.group(1)
                    val = mapping.get(key) or global_vars.get(key)
                    if val:
                        # remove text and add picture to slide
                        shape.text = ''
                        if isinstance(val, (bytes, bytearray)):
                            bio = io.BytesIO(val)
                            slide.shapes.add_picture(bio, Inches(1), Inches(1), width=Inches(4))
                        else:
                            try:
                                p = Path(val)
                                slide.shapes.add_picture(str(p), Inches(1), Inches(1), width=Inches(4))
                            except Exception:
                                pass
                        continue
                def repl(mo):
                    nm = mo.group(1)
                    if ':' in nm:
                        nm = nm.split(':',1)[1]
                    return str(mapping.get(nm, global_vars.get(nm, mo.group(0))))
                new_text = VAR_RE.sub(repl, text)
                if new_text != text:
                    shape.text = new_text

    def save(self, out_path: str):
        if self.prs is None:
            self.load()
        Path(out_path).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(out_path)

